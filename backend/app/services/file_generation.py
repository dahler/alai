"""
File Generation Service.
Generates downloadable files (DOCX, XLSX, PPTX, PDF, CSV) from structured content.
"""

import csv
import html
import io
import json
import os
import time
import uuid
from pathlib import Path
from typing import Optional


def _safe(text) -> str:
    """Escape XML special chars so reportlab Paragraph doesn't choke on them."""
    return html.escape(str(text))

# Absolute path anchored to this file's location so it never depends on CWD
GENERATED_DIR = Path(__file__).parent.parent.parent / "generated_files"
FILE_TTL_SECONDS = 3600  # files expire after 1 hour


def _content_type_for_ext(ext: str) -> str:
    return {
        ".csv": "text/csv",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".pdf": "application/pdf",
    }.get(ext, "application/octet-stream")


class FileGenerationService:
    def __init__(self):
        GENERATED_DIR.mkdir(parents=True, exist_ok=True)

    def _cleanup_old(self):
        now = time.time()
        for f in GENERATED_DIR.iterdir():
            if f.is_file() and (now - f.stat().st_mtime) > FILE_TTL_SECONDS:
                try:
                    f.unlink()
                except OSError:
                    pass

    def _save(self, data: bytes, original_filename: str, ext: str) -> dict:
        self._cleanup_old()
        file_id = str(uuid.uuid4())
        file_path = GENERATED_DIR / f"{file_id}{ext}"
        meta_path = GENERATED_DIR / f"{file_id}.meta.json"
        file_path.write_bytes(data)
        meta_path.write_text(json.dumps({
            "filename": original_filename,
            "content_type": _content_type_for_ext(ext),
            "ext": ext,
        }))
        return {
            "file_id": file_id,
            "filename": original_filename,
            "download_url": f"/api/files/download/{file_id}",
        }

    def get_file(self, file_id: str) -> Optional[tuple]:
        """Return (Path, original_filename, content_type) or None."""
        meta_path = GENERATED_DIR / f"{file_id}.meta.json"
        if not meta_path.exists():
            return None
        meta = json.loads(meta_path.read_text())
        file_path = GENERATED_DIR / f"{file_id}{meta['ext']}"
        if not file_path.exists():
            return None
        return file_path, meta["filename"], meta["content_type"]

    # ------------------------------------------------------------------
    # CSV
    # ------------------------------------------------------------------

    def generate_csv(self, filename: str, headers: list, rows: list) -> dict:
        buf = io.StringIO()
        writer = csv.writer(buf)
        if headers:
            writer.writerow(headers)
        writer.writerows(rows)
        data = buf.getvalue().encode("utf-8-sig")
        fname = filename if filename.endswith(".csv") else filename + ".csv"
        return self._save(data, fname, ".csv")

    # ------------------------------------------------------------------
    # Excel
    # ------------------------------------------------------------------

    def generate_excel(self, filename: str, sheets: list, title: str = "") -> dict:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment
        from openpyxl.utils import get_column_letter

        wb = Workbook()
        wb.remove(wb.active)  # remove default sheet

        for sheet_def in sheets:
            sheet_name = sheet_def.get("name", "Sheet1")[:31]
            ws = wb.create_sheet(title=sheet_name)

            headers = sheet_def.get("headers", [])
            rows = sheet_def.get("rows", [])

            # Write header row
            if headers:
                for col_idx, header in enumerate(headers, 1):
                    cell = ws.cell(row=1, column=col_idx, value=str(header))
                    cell.font = Font(bold=True, color="FFFFFF")
                    cell.fill = PatternFill("solid", fgColor="4472C4")
                    cell.alignment = Alignment(horizontal="center")

            # Write data rows
            for row_idx, row in enumerate(rows, 2 if headers else 1):
                for col_idx, value in enumerate(row, 1):
                    ws.cell(row=row_idx, column=col_idx, value=value)

            # Auto-width columns
            for col_idx in range(1, len(headers) + 1):
                col_letter = get_column_letter(col_idx)
                max_len = max(
                    (len(str(ws.cell(row=r, column=col_idx).value or "")) for r in range(1, ws.max_row + 1)),
                    default=10,
                )
                ws.column_dimensions[col_letter].width = min(max_len + 4, 60)

        buf = io.BytesIO()
        wb.save(buf)
        fname = filename if filename.endswith(".xlsx") else filename + ".xlsx"
        return self._save(buf.getvalue(), fname, ".xlsx")

    # ------------------------------------------------------------------
    # DOCX
    # ------------------------------------------------------------------

    def generate_docx(self, filename: str, title: str, sections: list) -> dict:
        from docx import Document
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()

        # Title
        if title:
            title_para = doc.add_heading(title, level=0)
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        for section in sections:
            heading = section.get("heading", "")
            level = section.get("level", 1)
            content = section.get("content", "")
            bullets = section.get("bullets", [])
            table_data = section.get("table")

            if heading:
                doc.add_heading(heading, level=min(level, 9))

            if content:
                doc.add_paragraph(content)

            for bullet in bullets:
                doc.add_paragraph(str(bullet), style="List Bullet")

            if table_data:
                headers = table_data.get("headers", [])
                rows = table_data.get("rows", [])
                if headers or rows:
                    num_cols = len(headers) if headers else (len(rows[0]) if rows else 1)
                    tbl = doc.add_table(rows=1 + len(rows), cols=num_cols)
                    tbl.style = "Table Grid"
                    if headers:
                        for i, h in enumerate(headers):
                            cell = tbl.rows[0].cells[i]
                            cell.text = str(h)
                            cell.paragraphs[0].runs[0].bold = True
                    for row_idx, row in enumerate(rows, 1):
                        for col_idx, val in enumerate(row):
                            tbl.rows[row_idx].cells[col_idx].text = str(val)
                    doc.add_paragraph()

        buf = io.BytesIO()
        doc.save(buf)
        fname = filename if filename.endswith(".docx") else filename + ".docx"
        return self._save(buf.getvalue(), fname, ".docx")

    # ------------------------------------------------------------------
    # PPTX
    # ------------------------------------------------------------------

    def generate_pptx(self, filename: str, title: str, slides: list) -> dict:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Title slide
        if title:
            layout = prs.slide_layouts[0]  # Title Slide
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title
            if slide.placeholders[1].has_text_frame:
                slide.placeholders[1].text = ""

        for slide_def in slides:
            slide_title = slide_def.get("title", "")
            content = slide_def.get("content", "")
            bullets = slide_def.get("bullets", [])
            table_data = slide_def.get("table")

            has_content = bool(content or bullets or table_data)
            layout = prs.slide_layouts[1] if has_content else prs.slide_layouts[5]
            slide = prs.slides.add_slide(layout)

            if slide.shapes.title:
                slide.shapes.title.text = slide_title

            if has_content and len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                if content:
                    tf.text = content
                for bullet in bullets:
                    p = tf.add_paragraph()
                    p.text = str(bullet)
                    p.level = 0

            # Add table if present (no content placeholder used)
            if table_data:
                headers = table_data.get("headers", [])
                rows = table_data.get("rows", [])
                num_rows = len(rows) + (1 if headers else 0)
                num_cols = len(headers) if headers else (len(rows[0]) if rows else 1)
                if num_rows > 0 and num_cols > 0:
                    tbl = slide.shapes.add_table(
                        num_rows, num_cols,
                        Inches(0.5), Inches(2.5),
                        slide_width - Inches(1), slide_height - Inches(3.5)
                    ).table
                    row_offset = 0
                    if headers:
                        for ci, h in enumerate(headers):
                            cell = tbl.cell(0, ci)
                            cell.text = str(h)
                            cell.text_frame.paragraphs[0].runs[0].font.bold = True
                        row_offset = 1
                    for ri, row in enumerate(rows):
                        for ci, val in enumerate(row):
                            tbl.cell(ri + row_offset, ci).text = str(val)

        buf = io.BytesIO()
        prs.save(buf)
        fname = filename if filename.endswith(".pptx") else filename + ".pptx"
        return self._save(buf.getvalue(), fname, ".pptx")

    # ------------------------------------------------------------------
    # PDF
    # ------------------------------------------------------------------

    def generate_pdf(self, filename: str, title: str, sections: list) -> dict:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.lib import colors
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, ListFlowable, ListItem
        )

        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4, leftMargin=2*cm, rightMargin=2*cm,
                                topMargin=2*cm, bottomMargin=2*cm)
        styles = getSampleStyleSheet()
        story = []

        heading_styles = [
            ParagraphStyle("H1", parent=styles["Heading1"], fontSize=18, spaceAfter=12),
            ParagraphStyle("H2", parent=styles["Heading2"], fontSize=14, spaceAfter=8),
            ParagraphStyle("H3", parent=styles["Heading3"], fontSize=12, spaceAfter=6),
        ]

        if title:
            story.append(Paragraph(_safe(title), heading_styles[0]))
            story.append(Spacer(1, 0.5*cm))

        for section in sections:
            heading = section.get("heading", "")
            level = max(1, min(section.get("level", 1), 3))
            content = section.get("content", "")
            bullets = section.get("bullets", [])
            table_data = section.get("table")

            if heading:
                story.append(Paragraph(_safe(heading), heading_styles[level - 1]))

            if content:
                story.append(Paragraph(_safe(content), styles["Normal"]))
                story.append(Spacer(1, 0.3*cm))

            if bullets:
                items = [ListItem(Paragraph(_safe(b), styles["Normal"])) for b in bullets]
                story.append(ListFlowable(items, bulletType="bullet"))
                story.append(Spacer(1, 0.3*cm))

            if table_data:
                headers = table_data.get("headers", [])
                rows = table_data.get("rows", [])
                table_rows = []
                if headers:
                    table_rows.append([_safe(h) for h in headers])
                for row in rows:
                    table_rows.append([_safe(v) for v in row])
                if table_rows:
                    tbl = Table(table_rows)
                    tbl.setStyle(TableStyle([
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4472C4")),
                        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#EEF3FB")]),
                        ("ALIGN", (0, 0), (-1, -1), "LEFT"),
                        ("PADDING", (0, 0), (-1, -1), 6),
                    ]))
                    story.append(tbl)
                    story.append(Spacer(1, 0.4*cm))

        doc.build(story)
        fname = filename if filename.endswith(".pdf") else filename + ".pdf"
        return self._save(buf.getvalue(), fname, ".pdf")

    # ------------------------------------------------------------------
    # Template-based fill (style-preserving)
    # ------------------------------------------------------------------

    def fill_docx_template(self, template_path: str, filename: str, title: str, sections: list) -> dict:
        """Open an existing DOCX template and insert AI content after matching headings."""
        from docx import Document
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
        from copy import deepcopy

        doc = Document(template_path)
        body = doc.element.body

        def _make_para_elem(text: str, style_name: str = "Normal"):
            """Create a <w:p> element with the given text and style."""
            p_elem = OxmlElement("w:p")
            pPr = OxmlElement("w:pPr")
            pStyle = OxmlElement("w:pStyle")
            # Map style name to docx style ID (spaces → no spaces, Word convention)
            style_id = style_name.replace(" ", "")
            pStyle.set(qn("w:val"), style_id)
            pPr.append(pStyle)
            p_elem.append(pPr)
            r = OxmlElement("w:r")
            t = OxmlElement("w:t")
            t.text = text
            t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
            r.append(t)
            p_elem.append(r)
            return p_elem

        # Build index: heading text (lower) → body child index
        body_children = list(body)
        heading_positions: list[tuple[str, int]] = []
        for i, child in enumerate(body_children):
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag == "p":
                pStyle = child.find(f".//{qn('w:pStyle')}")
                if pStyle is not None:
                    style_val = pStyle.get(qn("w:val"), "")
                    if style_val.startswith("Heading") or style_val.startswith("1") or style_val.startswith("2"):
                        pass  # handled below
                # Also check via para style name in document
            # Use python-docx paragraph objects for style checking
        para_map: dict[str, int] = {}
        body_para_indices: dict[int, int] = {}  # para index → body child index
        para_counter = 0
        for bi, child in enumerate(body_children):
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag == "p":
                if para_counter < len(doc.paragraphs):
                    para = doc.paragraphs[para_counter]
                    if para.style.name.startswith("Heading") and para.text.strip():
                        para_map[para.text.strip().lower()] = bi
                    body_para_indices[para_counter] = bi
                    para_counter += 1

        section_map = {s.get("heading", "").strip().lower(): s for s in sections if s.get("heading")}

        # Process in reverse body-position order to avoid index shifts
        for heading_lower, body_idx in sorted(para_map.items(), key=lambda x: x[1], reverse=True):
            sec = section_map.get(heading_lower)
            if not sec:
                continue

            insert_at = body_idx + 1
            content_text = sec.get("content", "")
            bullets = sec.get("bullets", [])
            table_data = sec.get("table")

            # Build elements in reverse so inserting at same position keeps order
            elems_to_insert = []
            if content_text:
                elems_to_insert.append(_make_para_elem(content_text, "Normal"))
            for b in bullets:
                elems_to_insert.append(_make_para_elem(str(b), "ListBullet"))

            if table_data:
                headers = table_data.get("headers", [])
                rows = table_data.get("rows", [])
                num_cols = len(headers) if headers else (len(rows[0]) if rows else 0)
                if num_cols:
                    tbl = doc.add_table(rows=1 + len(rows), cols=num_cols)
                    try:
                        tbl.style = "TableGrid"
                    except Exception:
                        pass
                    if headers:
                        for ci, h in enumerate(headers):
                            cell = tbl.rows[0].cells[ci]
                            cell.text = str(h)
                            if cell.paragraphs[0].runs:
                                cell.paragraphs[0].runs[0].bold = True
                    for ri, row in enumerate(rows, 1):
                        for ci, val in enumerate(row):
                            tbl.rows[ri].cells[ci].text = str(val)
                    tbl_elem = tbl._element
                    tbl_elem.getparent().remove(tbl_elem)
                    elems_to_insert.append(tbl_elem)

            for elem in reversed(elems_to_insert):
                body.insert(insert_at, elem)

        buf = io.BytesIO()
        doc.save(buf)
        fname = filename if filename.endswith(".docx") else filename + ".docx"
        return self._save(buf.getvalue(), fname, ".docx")

    def fill_xlsx_template(self, template_path: str, filename: str, sheets: list) -> dict:
        """Open an existing XLSX template and append data rows after the header rows."""
        import openpyxl

        wb = openpyxl.load_workbook(template_path)

        for sheet_def in sheets:
            sheet_name = sheet_def.get("name", "")
            rows = sheet_def.get("rows", [])
            if not rows:
                continue

            # Find matching sheet (case-insensitive)
            ws = None
            for sn in wb.sheetnames:
                if sn.lower() == sheet_name.lower():
                    ws = wb[sn]
                    break
            if ws is None:
                ws = wb.active  # fall back to first sheet

            # Find first empty row (after headers / existing data)
            start_row = ws.max_row + 1
            for row in rows:
                for ci, val in enumerate(row, 1):
                    ws.cell(row=start_row, column=ci, value=val)
                start_row += 1

        buf = io.BytesIO()
        wb.save(buf)
        fname = filename if filename.endswith(".xlsx") else filename + ".xlsx"
        return self._save(buf.getvalue(), fname, ".xlsx")

    def fill_pptx_template(self, template_path: str, filename: str, title: str, slides: list) -> dict:
        """Open an existing PPTX template and fill slide content matching by title."""
        from pptx import Presentation
        from pptx.util import Inches
        from copy import deepcopy

        prs = Presentation(template_path)

        # Build slide title → slide index map
        slide_map: dict[str, int] = {}
        for i, slide in enumerate(prs.slides):
            shape = slide.shapes.title
            if shape and shape.has_text_frame:
                slide_map[shape.text_frame.text.strip().lower()] = i

        for slide_def in slides:
            slide_title = slide_def.get("title", slide_def.get("heading", ""))
            content = slide_def.get("content", "")
            bullets = slide_def.get("bullets", [])
            table_data = slide_def.get("table")

            idx = slide_map.get(slide_title.strip().lower())
            if idx is None:
                # Append a new slide using the last layout used
                layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(layout)
                if slide.shapes.title:
                    slide.shapes.title.text = slide_title
            else:
                slide = prs.slides[idx]

            # Fill the body placeholder
            body_ph = None
            for ph in slide.placeholders:
                if ph.placeholder_format.idx != 0:  # skip title placeholder
                    body_ph = ph
                    break

            if body_ph and body_ph.has_text_frame:
                tf = body_ph.text_frame
                tf.clear()
                if content:
                    tf.text = content
                for b in bullets:
                    p = tf.add_paragraph()
                    p.text = str(b)

            if table_data:
                headers = table_data.get("headers", [])
                rows = table_data.get("rows", [])
                num_cols = len(headers) if headers else (len(rows[0]) if rows else 0)
                num_rows_total = len(rows) + (1 if headers else 0)
                if num_cols and num_rows_total:
                    tbl = slide.shapes.add_table(
                        num_rows_total, num_cols,
                        Inches(0.5), Inches(2.5),
                        prs.slide_width - Inches(1), prs.slide_height - Inches(3.5)
                    ).table
                    row_offset = 0
                    if headers:
                        for ci, h in enumerate(headers):
                            tbl.cell(0, ci).text = str(h)
                            tbl.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
                        row_offset = 1
                    for ri, row in enumerate(rows):
                        for ci, val in enumerate(row):
                            tbl.cell(ri + row_offset, ci).text = str(val)

        buf = io.BytesIO()
        prs.save(buf)
        fname = filename if filename.endswith(".pptx") else filename + ".pptx"
        return self._save(buf.getvalue(), fname, ".pptx")

    # ------------------------------------------------------------------
    # Unified dispatcher
    # ------------------------------------------------------------------

    def generate(self, fmt: str, filename: str, content: dict, template_file_path: Optional[str] = None) -> dict:
        """
        Dispatch to the right generator. If template_file_path is provided and the
        file exists, use the style-preserving fill methods instead of generating from scratch.
        """
        """
        Dispatch to the right generator.

        content schema (all formats share same structure):
          title: str
          sections: list of {heading, level, content, bullets, table: {headers, rows}}
          sheets:   list of {name, headers, rows}   — xlsx/csv only
        """
        title = content.get("title", filename)
        sections = content.get("sections", [])
        sheets = content.get("sheets", [])

        tpl = template_file_path if (template_file_path and Path(template_file_path).exists()) else None

        if fmt == "csv":
            # Use first sheet or first section's table
            if sheets:
                s = sheets[0]
                return self.generate_csv(filename, s.get("headers", []), s.get("rows", []))
            for sec in sections:
                if sec.get("table"):
                    t = sec["table"]
                    return self.generate_csv(filename, t.get("headers", []), t.get("rows", []))
            return self.generate_csv(filename, [], [])

        elif fmt == "xlsx":
            if not sheets:
                sheets = [
                    {
                        "name": sec.get("heading", "Sheet1")[:31] or "Sheet1",
                        "headers": sec["table"].get("headers", []),
                        "rows": sec["table"].get("rows", []),
                    }
                    for sec in sections if sec.get("table")
                ]
            if not sheets:
                sheets = [{"name": "Sheet1", "headers": [], "rows": []}]
            if tpl and tpl.endswith(".xlsx"):
                return self.fill_xlsx_template(tpl, filename, sheets)
            return self.generate_excel(filename, sheets, title)

        elif fmt == "docx":
            if tpl and tpl.endswith(".docx"):
                return self.fill_docx_template(tpl, filename, title, sections)
            return self.generate_docx(filename, title, sections)

        elif fmt == "pptx":
            slides = content.get("slides", sections)
            if tpl and tpl.endswith(".pptx"):
                return self.fill_pptx_template(tpl, filename, title, slides)
            return self.generate_pptx(filename, title, slides)

        elif fmt == "pdf":
            return self.generate_pdf(filename, title, sections)

        else:
            raise ValueError(f"Unsupported format: {fmt}")
