"""
Document extraction service for PDF, Word, PowerPoint, Excel and text files.
Extracts text content to send to the AI model.
"""

import time
import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


def log(message: str):
    timestamp = time.strftime("%H:%M:%S")
    print(f"[{timestamp}] [DOCUMENT] {message}")


class DocumentService:
    """Service for extracting text from documents."""

    SUPPORTED_EXTENSIONS = {
        '.pdf', '.docx', '.pptx', '.xlsx',
        '.txt', '.md', '.json', '.xml', '.html', '.css', '.js',
    }

    async def extract_text(self, file_path: str) -> Optional[str]:
        path = Path(file_path)
        log(f"Extracting from: {file_path}")
        log(f"Path exists: {path.exists()}")

        if not path.exists():
            log("File not found")
            return None

        extension = path.suffix.lower()
        log(f"File extension: {extension}")

        try:
            if extension == '.pdf':
                return await self._extract_pdf(path)
            elif extension == '.docx':
                return await self._extract_docx(path)
            elif extension == '.pptx':
                return await self._extract_pptx(path)
            elif extension == '.xlsx':
                return await self._extract_xlsx(path)
            elif extension in {
                '.txt', '.md', '.json', '.xml',
                '.html', '.css', '.js',
            }:
                return await self._extract_text_file(path)
            else:
                log(f"Unsupported file type: {extension}")
                return None
        except Exception as e:
            import traceback
            log(f"Extraction error: {e}")
            log(traceback.format_exc())
            return None

    async def _extract_pdf(self, path: Path) -> Optional[str]:
        log("PDF extraction starting...")
        start_time = time.time()

        try:
            import fitz  # PyMuPDF
            log(f"Using PyMuPDF v{fitz.version[0]}")

            text_parts = []
            with fitz.open(path) as doc:
                page_count = len(doc)
                log(f"PDF has {page_count} page(s)")
                for page_num, page in enumerate(doc, 1):
                    text = page.get_text()
                    if text.strip():
                        text_parts.append(
                            f"--- Page {page_num} ---\n{text}"
                        )

            if text_parts:
                full_text = "\n\n".join(text_parts)
                elapsed = time.time() - start_time
                log(
                    f"Extracted {len(full_text)} chars "
                    f"from {len(text_parts)} page(s) in {elapsed:.2f}s"
                )
                return full_text

            log("No text found in PDF (may be image-based)")
            return None

        except ImportError as e:
            log(f"PyMuPDF not installed: {e}, trying pypdf...")
            try:
                from pypdf import PdfReader
                reader = PdfReader(path)
                text_parts = []
                for page_num, page in enumerate(reader.pages, 1):
                    text = page.extract_text()
                    if text and text.strip():
                        text_parts.append(
                            f"--- Page {page_num} ---\n{text}"
                        )
                if text_parts:
                    full_text = "\n\n".join(text_parts)
                    elapsed = time.time() - start_time
                    log(
                        f"Extracted {len(full_text)} chars "
                        f"via pypdf in {elapsed:.2f}s"
                    )
                    return full_text
            except ImportError:
                log("No PDF library installed. Run: pip install pymupdf")
                return None
        except Exception as e:
            import traceback
            log(f"PDF extraction error: {e}")
            log(traceback.format_exc())
            return None

        return None

    async def _extract_docx(self, path: Path) -> Optional[str]:
        log(f"Extracting docx: {path.name}")
        try:
            from docx import Document
            doc = Document(path)
            parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(
                        c.text for c in row.cells if c.text.strip()
                    )
                    if row_text:
                        parts.append(row_text)
            full_text = "\n".join(parts)
            log(f"Extracted {len(full_text)} chars from docx")
            return full_text or None
        except ImportError:
            log("python-docx not installed. Run: pip install python-docx")
            return None

    async def _extract_pptx(self, path: Path) -> Optional[str]:
        log(f"Extracting pptx: {path.name}")
        try:
            from pptx import Presentation
            prs = Presentation(path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                slide_texts = [
                    shape.text
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ]
                if slide_texts:
                    parts.append(
                        f"--- Slide {i} ---\n" + "\n".join(slide_texts)
                    )
            full_text = "\n\n".join(parts)
            log(f"Extracted {len(full_text)} chars from pptx")
            return full_text or None
        except ImportError:
            log("python-pptx not installed. Run: pip install python-pptx")
            return None

    async def _extract_xlsx(self, path: Path) -> Optional[str]:
        log(f"Extracting xlsx: {path.name}")
        try:
            import openpyxl
            wb = openpyxl.load_workbook(
                path, read_only=True, data_only=True
            )
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                parts.append(f"--- Sheet: {sheet_name} ---")
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join(
                        str(v) for v in row if v is not None
                    )
                    if row_text.strip():
                        parts.append(row_text)
            wb.close()
            full_text = "\n".join(parts)
            log(f"Extracted {len(full_text)} chars from xlsx")
            return full_text or None
        except ImportError:
            log("openpyxl not installed. Run: pip install openpyxl")
            return None

    async def _extract_text_file(self, path: Path) -> Optional[str]:
        log(f"Extracting text file: {path.name}")
        try:
            text = path.read_text(encoding='utf-8')
            log(f"Extracted {len(text)} chars (UTF-8)")
            return text
        except UnicodeDecodeError:
            try:
                text = path.read_text(encoding='latin-1')
                log(f"Extracted {len(text)} chars (latin-1 fallback)")
                return text
            except Exception as e:
                log(f"Text extraction error: {e}")
                return None

    def truncate_text(self, text: str, max_chars: int = 15000) -> str:
        if len(text) <= max_chars:
            return text

        truncated = text[:max_chars]
        last_period = truncated.rfind('.')
        last_newline = truncated.rfind('\n')
        break_point = max(last_period, last_newline)

        if break_point > max_chars * 0.8:
            truncated = truncated[:break_point + 1]

        return truncated + "\n\n[Document truncated due to length...]"
