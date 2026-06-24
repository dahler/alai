"""
CRUD endpoints for report templates.

Access rules:
- Any authenticated user can create/read/edit/delete their own templates.
- Admin users can additionally publish templates as company-wide and manage all templates.
- GET /templates returns the calling user's own templates + all company-wide templates.
"""

import io
import json
import uuid
from pathlib import Path
from typing import List, Optional

from fastapi import APIRouter, Depends, File, HTTPException, UploadFile, status
from pydantic import BaseModel
from sqlalchemy import or_, select
from sqlalchemy.ext.asyncio import AsyncSession

from app.database import get_db
from app.middleware.auth import get_current_user
from app.models.report_template import ReportTemplate
from app.models.user import User

# Directory for permanent template files
TEMPLATES_DIR = Path(__file__).parent.parent.parent / "template_files"
TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)

# Directory for temporary files (pending template creation)
TEMPLATES_TEMP_DIR = Path(__file__).parent.parent.parent / "template_files" / "tmp"
TEMPLATES_TEMP_DIR.mkdir(parents=True, exist_ok=True)

router = APIRouter(prefix="/templates", tags=["templates"])


# ---------------------------------------------------------------------------
# Schemas
# ---------------------------------------------------------------------------

class TemplateSectionIn(BaseModel):
    heading: str
    level: int = 1
    placeholder: str = ""
    has_table: bool = False
    table_headers: List[str] = []


class TemplateCreate(BaseModel):
    name: str
    description: Optional[str] = None
    format: str  # pdf | docx | xlsx | pptx
    sections: List[TemplateSectionIn]
    keywords: Optional[str] = None
    temp_file_id: Optional[str] = None  # returned by /extract-headings


class TemplateUpdate(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None
    format: Optional[str] = None
    sections: Optional[List[TemplateSectionIn]] = None
    keywords: Optional[str] = None
    temp_file_id: Optional[str] = None


class TemplateSectionOut(BaseModel):
    heading: str
    level: int
    placeholder: str
    has_table: bool
    table_headers: List[str]


class TemplateOut(BaseModel):
    id: int
    name: str
    description: Optional[str]
    format: str
    sections: List[TemplateSectionOut]
    keywords: Optional[str]
    owner_id: Optional[int]
    is_company_wide: bool
    is_mine: bool
    has_template_file: bool  # whether a source file is stored for style-preserving generation


def _to_out(tpl: ReportTemplate, current_user_id: int) -> TemplateOut:
    sections = json.loads(tpl.sections_json) if tpl.sections_json else []
    return TemplateOut(
        id=tpl.id,
        name=tpl.name,
        description=tpl.description,
        format=tpl.format,
        sections=[TemplateSectionOut(**s) for s in sections],
        keywords=tpl.keywords,
        owner_id=tpl.owner_id,
        is_company_wide=tpl.is_company_wide,
        is_mine=(tpl.owner_id == current_user_id),
        has_template_file=bool(tpl.template_file_path and Path(tpl.template_file_path).exists()),
    )


def _commit_temp_file(temp_file_id: str, template_id: int, fmt: str) -> Optional[str]:
    """Move a temp file to permanent storage. Returns the permanent path or None."""
    allowed_ext = {"docx": ".docx", "xlsx": ".xlsx", "pptx": ".pptx", "pdf": ".pdf"}
    ext = allowed_ext.get(fmt)
    if not ext or not temp_file_id:
        return None
    # Try all possible extensions for the temp file
    for try_ext in [ext, ".docx", ".xlsx", ".pptx", ".pdf"]:
        src = TEMPLATES_TEMP_DIR / f"{temp_file_id}{try_ext}"
        if src.exists():
            dest = TEMPLATES_DIR / f"tpl_{template_id}{try_ext}"
            src.rename(dest)
            return str(dest)
    return None


def _validate_format(fmt: str) -> str:
    allowed = {"pdf", "docx", "xlsx", "pptx"}
    if fmt not in allowed:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=f"format must be one of {sorted(allowed)}",
        )
    return fmt


# ---------------------------------------------------------------------------
# Endpoints
# ---------------------------------------------------------------------------

@router.get("", response_model=List[TemplateOut])
async def list_templates(
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Return templates visible to the current user: own + company-wide."""
    result = await db.execute(
        select(ReportTemplate).where(
            or_(
                ReportTemplate.owner_id == current_user.id,
                ReportTemplate.is_company_wide == True,  # noqa: E712
            )
        ).order_by(ReportTemplate.is_company_wide, ReportTemplate.name)
    )
    templates = result.scalars().all()
    return [_to_out(t, current_user.id) for t in templates]


@router.post("", response_model=TemplateOut, status_code=status.HTTP_201_CREATED)
async def create_template(
    body: TemplateCreate,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    _validate_format(body.format)
    sections_data = [s.model_dump() for s in body.sections]
    tpl = ReportTemplate(
        name=body.name,
        description=body.description,
        format=body.format,
        sections_json=json.dumps(sections_data),
        keywords=body.keywords,
        owner_id=current_user.id,
        is_company_wide=False,
    )
    db.add(tpl)
    await db.commit()
    await db.refresh(tpl)
    if body.temp_file_id:
        perm_path = _commit_temp_file(body.temp_file_id, tpl.id, body.format)
        if perm_path:
            tpl.template_file_path = perm_path
            await db.commit()
    return _to_out(tpl, current_user.id)


@router.get("/{template_id}", response_model=TemplateOut)
async def get_template(
    template_id: int,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    tpl = await _get_accessible(template_id, current_user, db)
    return _to_out(tpl, current_user.id)


@router.put("/{template_id}", response_model=TemplateOut)
async def update_template(
    template_id: int,
    body: TemplateUpdate,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    tpl = await _get_owned(template_id, current_user, db)
    if body.name is not None:
        tpl.name = body.name
    if body.description is not None:
        tpl.description = body.description
    if body.format is not None:
        _validate_format(body.format)
        tpl.format = body.format
    if body.sections is not None:
        tpl.sections_json = json.dumps([s.model_dump() for s in body.sections])
    if body.keywords is not None:
        tpl.keywords = body.keywords
    if body.temp_file_id:
        fmt = body.format or tpl.format
        perm_path = _commit_temp_file(body.temp_file_id, tpl.id, fmt)
        if perm_path:
            # Remove old file if it exists
            if tpl.template_file_path:
                try:
                    Path(tpl.template_file_path).unlink(missing_ok=True)
                except OSError:
                    pass
            tpl.template_file_path = perm_path
    await db.commit()
    await db.refresh(tpl)
    return _to_out(tpl, current_user.id)


@router.delete("/{template_id}", status_code=status.HTTP_204_NO_CONTENT)
async def delete_template(
    template_id: int,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    tpl = await _get_owned(template_id, current_user, db)
    await db.delete(tpl)
    await db.commit()


@router.patch("/{template_id}/company-wide", response_model=TemplateOut)
async def toggle_company_wide(
    template_id: int,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Toggle company-wide visibility. Any user can do this for their own templates."""
    tpl = await _get_owned(template_id, current_user, db)
    tpl.is_company_wide = not tpl.is_company_wide
    await db.commit()
    await db.refresh(tpl)
    return _to_out(tpl, current_user.id)


# ---------------------------------------------------------------------------
# Heading extraction
# ---------------------------------------------------------------------------

class ExtractedHeading(BaseModel):
    heading: str
    level: int


class ExtractHeadingsResponse(BaseModel):
    headings: List[ExtractedHeading]
    temp_file_id: Optional[str] = None  # set only for DOCX/XLSX/PPTX (not PDF)


@router.post("/extract-headings", response_model=ExtractHeadingsResponse)
async def extract_headings(
    file: UploadFile = File(...),
    current_user: User = Depends(get_current_user),
):
    """
    Upload a DOCX or PDF file and return headings + a temp_file_id.
    The temp_file_id references the saved file so it can be attached to
    the template on creation, enabling style-preserving generation.
    """
    content = await file.read()
    name = (file.filename or "").lower()

    if name.endswith(".docx"):
        headings = _extract_docx_headings(content)
        # Save file for later use as generation template
        temp_id = str(uuid.uuid4())
        (TEMPLATES_TEMP_DIR / f"{temp_id}.docx").write_bytes(content)
        return ExtractHeadingsResponse(headings=headings, temp_file_id=temp_id)
    elif name.endswith(".xlsx"):
        # Extract sheet/column names as pseudo-headings
        headings = _extract_xlsx_headings(content)
        temp_id = str(uuid.uuid4())
        (TEMPLATES_TEMP_DIR / f"{temp_id}.xlsx").write_bytes(content)
        return ExtractHeadingsResponse(headings=headings, temp_file_id=temp_id)
    elif name.endswith(".pptx"):
        headings = _extract_pptx_headings(content)
        temp_id = str(uuid.uuid4())
        (TEMPLATES_TEMP_DIR / f"{temp_id}.pptx").write_bytes(content)
        return ExtractHeadingsResponse(headings=headings, temp_file_id=temp_id)
    elif name.endswith(".pdf"):
        headings = _extract_pdf_headings(content)
        # Save PDF so the vision model can use it as a visual style reference
        temp_id = str(uuid.uuid4())
        (TEMPLATES_TEMP_DIR / f"{temp_id}.pdf").write_bytes(content)
        return ExtractHeadingsResponse(headings=headings, temp_file_id=temp_id)
    else:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail="Only .docx, .xlsx, .pptx and .pdf files are supported.",
        )


def _extract_docx_headings(content: bytes) -> List[ExtractedHeading]:
    try:
        from docx import Document  # type: ignore
    except ImportError:
        raise HTTPException(status_code=500, detail="python-docx is not installed.")

    doc = Document(io.BytesIO(content))
    headings: List[ExtractedHeading] = []
    for para in doc.paragraphs:
        style = para.style.name  # e.g. "Heading 1", "Heading 2"
        if style.startswith("Heading") and para.text.strip():
            parts = style.split()
            level = int(parts[-1]) if parts and parts[-1].isdigit() else 1
            headings.append(ExtractedHeading(heading=para.text.strip(), level=min(level, 3)))

    if not headings:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail="No headings found. Make sure the document uses Word Heading styles (Heading 1, Heading 2, …).",
        )
    return headings


def _extract_xlsx_headings(content: bytes) -> List[ExtractedHeading]:
    try:
        import openpyxl  # type: ignore
    except ImportError:
        raise HTTPException(status_code=500, detail="openpyxl is not installed.")

    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    headings: List[ExtractedHeading] = []
    for sheet_name in wb.sheetnames:
        headings.append(ExtractedHeading(heading=sheet_name, level=1))
        ws = wb[sheet_name]
        first_row = next(ws.iter_rows(min_row=1, max_row=1, values_only=True), None)
        if first_row:
            for cell_value in first_row:
                if cell_value and isinstance(cell_value, str) and cell_value.strip():
                    headings.append(ExtractedHeading(heading=cell_value.strip(), level=2))
    wb.close()
    if not headings:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail="No sheets or column headers found in the spreadsheet.",
        )
    return headings


def _extract_pptx_headings(content: bytes) -> List[ExtractedHeading]:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        raise HTTPException(status_code=500, detail="python-pptx is not installed.")

    prs = Presentation(io.BytesIO(content))
    headings: List[ExtractedHeading] = []
    for i, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        if title_shape and title_shape.has_text_frame:
            text = title_shape.text_frame.text.strip()
            if text:
                headings.append(ExtractedHeading(heading=text, level=1))
                continue
        headings.append(ExtractedHeading(heading=f"Slide {i}", level=1))
    if not headings:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail="No slides found in the presentation.",
        )
    return headings


def _extract_pdf_headings(content: bytes) -> List[ExtractedHeading]:
    try:
        import fitz  # PyMuPDF  # type: ignore
    except ImportError:
        raise HTTPException(status_code=500, detail="PyMuPDF is not installed.")

    doc = fitz.open(stream=content, filetype="pdf")

    # Collect every text span with its font size
    spans = []
    for page in doc:
        for block in page.get_text("dict")["blocks"]:
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    text = span.get("text", "").strip()
                    size = round(span.get("size", 0), 1)
                    if text and len(text) < 300:
                        spans.append({"text": text, "size": size})

    if not spans:
        raise HTTPException(status_code=422, detail="No text found in PDF.")

    # Body text is the most common font size; headings are larger
    from statistics import median
    sizes = [s["size"] for s in spans]
    body_size = median(sizes)
    threshold = body_size + 1.5  # anything 1.5pt larger than median is a heading

    heading_spans = [s for s in spans if s["size"] > threshold]
    if not heading_spans:
        raise HTTPException(
            status_code=422,
            detail="Could not detect headings. PDF may not have distinct heading font sizes.",
        )

    max_size = max(s["size"] for s in heading_spans)

    headings: List[ExtractedHeading] = []
    seen: set = set()
    for s in heading_spans:
        text = s["text"]
        if text in seen:
            continue
        seen.add(text)
        # Map font size → heading level
        if s["size"] >= max_size - 0.5:
            level = 1
        elif s["size"] >= max_size - 3:
            level = 2
        else:
            level = 3
        headings.append(ExtractedHeading(heading=text, level=level))

    return headings


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

async def _get_accessible(
    template_id: int, current_user: User, db: AsyncSession
) -> ReportTemplate:
    result = await db.execute(
        select(ReportTemplate).where(
            ReportTemplate.id == template_id,
            or_(
                ReportTemplate.owner_id == current_user.id,
                ReportTemplate.is_company_wide == True,  # noqa: E712
            ),
        )
    )
    tpl = result.scalar_one_or_none()
    if not tpl:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Template not found")
    return tpl


async def _get_owned(
    template_id: int, current_user: User, db: AsyncSession
) -> ReportTemplate:
    result = await db.execute(
        select(ReportTemplate).where(ReportTemplate.id == template_id)
    )
    tpl = result.scalar_one_or_none()
    if not tpl:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Template not found",
        )
    if tpl.owner_id != current_user.id:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You can only modify your own templates.",
        )
    return tpl
